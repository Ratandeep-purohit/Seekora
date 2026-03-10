from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── COLOR PALETTE ────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x0F, 0x1A)   # deep navy-black
BG_CARD      = RGBColor(0x13, 0x15, 0x28)   # slightly lighter surface
ACCENT       = RGBColor(0x63, 0x66, 0xF1)   # indigo #6366f1
ACCENT_LIGHT = RGBColor(0xA5, 0xB4, 0xFC)   # soft lavender
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY         = RGBColor(0x94, 0xA3, 0xB8)
GREEN        = RGBColor(0x10, 0xB9, 0x81)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ─── HELPERS ──────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def set_bg(slide, color=BG_DARK):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)

def add_accent_bar(slide, y=Inches(0.55), height=Pt(3)):
    add_rect(slide, Inches(0.7), int(y), int(SLIDE_W - Inches(1.4)), int(height), ACCENT)

def add_slide_header(slide, title, subtitle=None):
    set_bg(slide)
    # Top accent line
    add_rect(slide, Inches(0.7), Inches(0.45), Inches(1.5), Pt(3), ACCENT)
    add_text(slide, title, Inches(0.7), Inches(0.5), Inches(11), Inches(0.8),
             size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.7), Inches(1.2), Inches(11), Inches(0.5),
                 size=16, color=GRAY)
    # bottom line
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), ACCENT)


def bullet_slide(slide, title, bullets, subtitle=None, columns=1):
    add_slide_header(slide, title, subtitle)
    y_start = Inches(1.9)
    avail_h  = Inches(5.2)
    col_w = (SLIDE_W - Inches(1.4)) / columns

    if columns == 1:
        for i, (icon, heading, body) in enumerate(bullets):
            card_y = y_start + i * Inches(0.85)
            add_rect(slide, Inches(0.7), card_y, Inches(0.05), Inches(0.6), ACCENT)
            add_text(slide, f"{icon}  {heading}", Inches(0.85), card_y, col_w - Inches(0.3),
                     Inches(0.4), size=16, bold=True, color=ACCENT_LIGHT)
            if body:
                add_text(slide, body, Inches(0.85), card_y + Inches(0.38),
                         col_w - Inches(0.3), Inches(0.4), size=13, color=GRAY)
    else:
        # 2-column grid
        for i, (icon, heading, body) in enumerate(bullets):
            col = i % columns
            row = i // columns
            cx = Inches(0.7) + col * col_w
            cy = y_start + row * Inches(1.4)
            # card background
            add_rect(slide, cx, cy, col_w - Inches(0.2), Inches(1.3), BG_CARD)
            add_rect(slide, cx, cy, Inches(0.05), Inches(1.3), ACCENT)
            add_text(slide, f"{icon}  {heading}", cx + Inches(0.15), cy + Inches(0.1),
                     col_w - Inches(0.4), Inches(0.45), size=15, bold=True, color=ACCENT_LIGHT)
            if body:
                add_text(slide, body, cx + Inches(0.15), cy + Inches(0.5),
                         col_w - Inches(0.4), Inches(0.7), size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1)

# Glow orb (faux — large semi-transparent circle)
add_rect(s1, Inches(8.5), Inches(-1), Inches(5), Inches(5), RGBColor(0x23, 0x27, 0x6A))

# Seekora wordmark
add_text(s1, "Seekora", Inches(1), Inches(1.8), Inches(8), Inches(2),
         size=72, bold=True, color=ACCENT_LIGHT, align=PP_ALIGN.LEFT)

# Tagline
add_text(s1, "A Next-Gen Search Engine — Built from Scratch",
         Inches(1), Inches(3.5), Inches(9), Inches(0.8),
         size=22, color=GRAY, align=PP_ALIGN.LEFT)

# Sub info
add_text(s1, "Second Year Project  |  Master of Computer Applications  |  2026",
         Inches(1), Inches(4.3), Inches(9), Inches(0.5),
         size=14, color=RGBColor(0x4F, 0x46, 0xE5), align=PP_ALIGN.LEFT)

# Bottom accent bar
add_rect(s1, Inches(0), SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), ACCENT)

# Accent line beside title
add_rect(s1, Inches(1), Inches(1.7), Inches(2), Pt(3), ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
bullet_slide(s2,
    "Problem Statement",
    [
        ("🔒", "Privacy Invasion",       "Google tracks every search, click, and location to serve targeted ads. Users have no control."),
        ("📢", "Ad Overload",            "Top search results are paid ads, not the most relevant content. Real results are buried below."),
        ("🔒", "Opaque Algorithms",      "Google's ranking algorithm is a black box. Bias and SEO manipulation go unchecked."),
        ("📦", "No Custom Control",      "Developers / organizations cannot customize search behavior for their own domain or use case."),
    ],
    subtitle="Why do we need a new search engine?"
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IS SEEKORA
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
bullet_slide(s3,
    "What is Seekora?",
    [
        ("🌐", "Full-Stack Search Engine",   "End-to-end system: crawler → indexer → ranker → frontend."),
        ("🕷️", "Own Web Crawler",            "Crawls real websites, respects robots.txt, stores structured data."),
        ("🧠", "NLP-Powered Queries",         "Tokenization, stemming, and spell correction for smarter matching."),
        ("🔗", "Federated Search",            "Aggregates Google CSE, Brave Search, Bing in parallel for maximum coverage."),
        ("🎨", "Premium Modern UI",           "Dark glassmorphism theme with 14 color themes, smooth Framer Motion animations."),
        ("🔐", "Privacy-First",               "Zero user tracking. No ads. No data sold."),
    ],
    subtitle="A complete, custom-built search platform"
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_slide_header(s4, "System Architecture", "How all the components connect")

# Layer diagram as colored boxes
def arch_box(slide, label, sublabel, x, y, w=Inches(3.6), h=Inches(0.85), color=ACCENT):
    add_rect(slide, x, y, w, h, BG_CARD)
    add_rect(slide, x, y, Inches(0.05), h, color)
    add_text(slide, label, x+Inches(0.12), y+Inches(0.05), w-Inches(0.2), Inches(0.4),
             size=14, bold=True, color=WHITE)
    add_text(slide, sublabel, x+Inches(0.12), y+Inches(0.45), w-Inches(0.2), Inches(0.35),
             size=11, color=GRAY)

cy = Inches(1.85)
arch_box(s4, "🌐  User Query (React Frontend)", "TypeScript · Zustand · TanStack Query · Framer Motion",
         Inches(0.6), cy, SLIDE_W - Inches(1.2), h=Inches(0.8), color=ACCENT)

cy += Inches(1.0)
arch_box(s4, "⚙️  Django REST API", "Receives query, dispatches to pipelines, merges results",
         Inches(0.6), cy, SLIDE_W - Inches(1.2), h=Inches(0.8), color=RGBColor(0x7C,0x3A,0xED))

cy += Inches(1.0)
arch_box(s4, "🕷️  Web Crawler", "BFS crawler · robots.txt · link discovery",
         Inches(0.5), cy, Inches(3.6), color=GREEN)
arch_box(s4, "🧠  NLP Processor", "NLTK · Stemming · Spell check",
         Inches(4.3), cy, Inches(3.6), color=ACCENT)
arch_box(s4, "🔗  Federated APIs", "Google CSE · Brave · Bing",
         Inches(8.1), cy, Inches(4.5), color=RGBColor(0xF5,0x9E,0x0B))

cy += Inches(1.0)
arch_box(s4, "🗃️  MySQL Database", "WebPage · SearchIndex · Image · Video models  —  Inverted Index with TF-IDF weights",
         Inches(0.6), cy, SLIDE_W - Inches(1.2), h=Inches(0.75), color=RGBColor(0x0D,0x94,0x88))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — WEB CRAWLER
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
bullet_slide(s5,
    "Web Crawler",
    [
        ("🚀", "BFS Link Discovery",       "Starts from seed URLs, follows all valid internal/external links recursively."),
        ("🤖", "robots.txt Compliance",    "Reads and respects crawl delay and disallow rules for every domain."),
        ("⚡", "Concurrent Fetching",      "ThreadPoolExecutor for parallel URL fetching — crawls faster."),
        ("📥", "Content Extraction",       "BeautifulSoup parses    title, description, body text, images, and videos."),
        ("🗄️", "Stores to MySQL",          "Saves structured WebPage objects with full-text content for indexing."),
    ],
    subtitle="Custom crawler engine built with Python & BeautifulSoup"
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — NLP & INDEXING
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
bullet_slide(s6,
    "NLP Query Processing & Indexing",
    [
        ("✂️", "Tokenization",             "Breaks raw query into individual words/tokens."),
        ("🌿", "Stemming (Porter)",        "Reduces words to root form — 'searching' → 'search'."),
        ("📝", "Spell Correction",         "Auto-suggests corrected query when typos are detected."),
        ("📊", "TF-IDF Ranking",           "Each word gets a relevance weight. Pages are ranked by total weighted score."),
        ("🔎", "Inverted Index",           "Word → [page_id, weight] map stored in MySQL SearchIndex table for O(1) lookup."),
    ],
    subtitle="Making searches relevant, not just keyword-matched"
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SEARCH VERTICALS
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
bullet_slide(s7,
    "Multi-Vertical Search",
    [
        ("🌐", "Web Search",   "Google CSE primary · Brave Search fallback · Local DB (TF-IDF) last resort"),
        ("🖼️", "Image Search", "Google Image API · Bing Image scraper fallback · 50 results per page"),
        ("📰", "News Search",  "Google News RSS feed · Real-time articles with source & timestamp"),
        ("🎥", "Video Search", "YouTube-focused scraping via DuckDuckGo · Direct YouTube fallback"),
    ],
    subtitle="Dedicated pipelines for each content type",
    columns=2
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
add_slide_header(s8, "Technology Stack")

techs = [
    ("Backend",  "Python · Django REST Framework · NLTK · BeautifulSoup · MySQL"),
    ("Frontend", "React · TypeScript · TailwindCSS · Framer Motion · Zustand"),
    ("Search",   "Google Custom Search · Brave Search · Bing · YouTube RSS"),
    ("DevOps",   "Git · GitHub · .env secrets · Django StatReloader"),
    ("Database", "MySQL · Django ORM · Paginator · Inverted Index"),
    ("UI/UX",    "Glassmorphism · Dark Mode · 14 Color Themes · Lucide Icons"),
]
cy = Inches(1.85)
col_w = (SLIDE_W - Inches(1.4)) / 2
for i, (label, desc) in enumerate(techs):
    col = i % 2
    row = i // 2
    cx = Inches(0.7) + col * col_w
    card_y = cy + row * Inches(1.3)
    add_rect(s8, cx, card_y, col_w - Inches(0.2), Inches(1.2), BG_CARD)
    add_rect(s8, cx, card_y, Inches(0.05), Inches(1.2), ACCENT)
    add_text(s8, label, cx + Inches(0.15), card_y + Inches(0.1),
             col_w - Inches(0.4), Inches(0.4), size=16, bold=True, color=ACCENT_LIGHT)
    add_text(s8, desc, cx + Inches(0.15), card_y + Inches(0.5),
             col_w - Inches(0.4), Inches(0.6), size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — KEY FEATURES
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
bullet_slide(s9,
    "Key Features",
    [
        ("⚡", "Live Crawling",          "Triggers real-time crawl if local index has < 10 results for a query."),
        ("🎨", "14 Color Themes",        "Full UI color scheme switches instantly — accent, background, glass cards all change."),
        ("🕐", "Live Clock Widget",       "Real-time date & time shown in the header of the home page."),
        ("🔍", "Autocomplete",           "Backend-powered suggestions from indexed keywords as you type."),
        ("📖", "Knowledge Panel",        "Rich summary card extracted from top result shown beside web results."),
        ("📄", "Pagination",             "Google-style page navigation for all search verticals."),
    ],
    subtitle="Beyond just search results",
    columns=2
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TEAM
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
add_slide_header(s10, "Our Team", "6 roles · 1 goal")

members = [
    ("Ratandeep",  "Backend & Web Crawler"),
    ("Jay",        "Database & Video Search"),
    ("Shivam",     "Search Ranking & Image Search"),
    ("Manav",      "News Pipeline & QA Testing"),
    ("Vanshika",   "Frontend UI/UX & Integration"),
    ("Jainika",    "Documentation & DevOps"),
]
col_w = (SLIDE_W - Inches(1.4)) / 3
cy = Inches(2.0)
for i, (name, role) in enumerate(members):
    col = i % 3
    row = i // 3
    cx = Inches(0.7) + col * col_w
    card_y = cy + row * Inches(1.8)
    add_rect(s10, cx, card_y, col_w - Inches(0.2), Inches(1.6), BG_CARD)
    # Avatar circle (faux)
    add_rect(s10, cx + Inches(0.15), card_y + Inches(0.15), Inches(0.55), Inches(0.55), ACCENT)
    add_text(s10, name[0], cx + Inches(0.25), card_y + Inches(0.17), Inches(0.4), Inches(0.4),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s10, name, cx + Inches(0.8), card_y + Inches(0.15), col_w - Inches(1.1),
             Inches(0.4), size=14, bold=True, color=WHITE)
    add_text(s10, role, cx + Inches(0.8), card_y + Inches(0.55), col_w - Inches(1.1),
             Inches(0.5), size=11, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank_layout)
set_bg(s11)
add_rect(s11, Inches(0), SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), ACCENT)

add_text(s11, "Key Takeaways", Inches(1), Inches(0.8), Inches(11), Inches(0.7),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

takeaways = [
    "A search box is NOT simple — it's a distributed system.",
    "NLP + Indexing is what separates search from ctrl+F.",
    "Federated search gives resilience when one source fails.",
    "Privacy-first design is possible without sacrificing quality.",
]
for i, t in enumerate(takeaways):
    cy = Inches(1.9) + i * Inches(0.9)
    add_rect(s11, Inches(1), cy, Inches(0.05), Inches(0.6), ACCENT)
    add_text(s11, t, Inches(1.2), cy + Inches(0.1), Inches(10.5), Inches(0.6),
             size=18, color=WHITE)

add_text(s11, "github.com/Ratandeep-purohit/Seekora",
         Inches(1), Inches(6.0), Inches(11), Inches(0.5),
         size=15, color=ACCENT_LIGHT, align=PP_ALIGN.CENTER)

add_text(s11, "Thank You  🙏",
         Inches(1), Inches(6.6), Inches(11), Inches(0.6),
         size=22, bold=True, color=ACCENT_LIGHT, align=PP_ALIGN.CENTER)


# ─── SAVE ─────────────────────────────────────────────────────────────────────
out_path = "Seekora_Presentation.pptx"
prs.save(out_path)
print(f"✅ Saved: {out_path}")
